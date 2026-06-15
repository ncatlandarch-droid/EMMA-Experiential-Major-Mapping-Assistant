"""
EMMA PSSG Pitch Deck — NC A&T Branded
For 2:45pm call with Dr. Williams, June 15, 2026
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import os

TEMPLATE = r'C:\Users\Chris\Downloads\NCATSTAT_Template.pptx'
OUTPUT   = r'C:\Users\Chris\Downloads\EMMA_PSSG_Pitch_Deck.pptx'

prs = Presentation(TEMPLATE)

# Delete all example slides from the template
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Layout references
L_TITLE      = prs.slide_layouts[0]
L_TRANSITION = prs.slide_layouts[2]
L_BULLET1    = prs.slide_layouts[5]
L_BULLET2    = prs.slide_layouts[6]
L_BULLET3    = prs.slide_layouts[7]
L_BULLET4    = prs.slide_layouts[8]

def add_bullets(slide, title, subtitle, bullets, font_size=Pt(13)):
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = title
        elif idx == 10:
            ph.text = subtitle
        elif idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, b in enumerate(bullets):
                if i == 0:
                    tf.paragraphs[0].text = b
                else:
                    p = tf.add_paragraph()
                    p.text = b
                for run in tf.paragraphs[i].runs:
                    run.font.size = font_size

# ─── SLIDE 1: TITLE ───
slide = prs.slides.add_slide(L_TITLE)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "EMMA: AI-Powered Student Success"

# ─── SLIDE 2: THE OPPORTUNITY ───
slide = prs.slides.add_slide(L_BULLET1)
add_bullets(slide,
    "The Opportunity",
    "FIPSE Postsecondary Student Success Grants (PSSG)",
    [
        "$45M total | ~6 early phase awards | ~$3.75M each",
        "48-month project period",
        "Deadline: June 29, 2026",
        "10% match (in-kind allowable)",
        "Requires: 1 evidence priority + 1 content priority",
        "10 bonus points for state-level partnership",
    ]
)

# ─── SLIDE 3: WHY NC A&T WINS ───
slide = prs.slides.add_slide(L_BULLET2)
add_bullets(slide,
    "Why NC A&T Wins This",
    "Competitive Positioning",
    [
        "Largest HBCU in the nation (~14,000 students)",
        "58% Pell Grant recipients (strongest tiebreaker)",
        "1890 Land-Grant institution",
        "EMMA already operational across 93 degree programs",
        "AI-powered coaching already built (Google Gemini)",
        "Dr. Williams' title mirrors the grant name",
    ]
)

# ─── SLIDE 4: WHAT IS EMMA ───
slide = prs.slides.add_slide(L_BULLET3)
add_bullets(slide,
    "What Is EMMA?",
    "Experiential Major Mapping Assistant",
    [
        "AI-powered platform mapping every student's 4-year journey",
        "4 phases: Explore > Engage > Develop > Launch",
        "Google Gemini AI coaching + neural voice",
        "Live BLS salary/employment data integration",
        "ISLA companion: licensure & accreditation tracking",
        "White-label architecture: 93 programs, all 8 colleges",
        "Working today — not a concept",
    ]
)

# ─── SLIDE 5: WHAT THE GRANT FUNDS ───
slide = prs.slides.add_slide(L_BULLET4)
add_bullets(slide,
    "What the Grant Funds",
    "A rigorous trial of what's already built",
    [
        "Independent evaluation (WWC standards)",
        "Project Director + student success coaches",
        "Technology scaling & infrastructure",
        "Student support services",
        "Years 1-2: CAES pilot (17 programs, ~3,000 students)",
        "Years 3-4: Scale to 4+ colleges (~6,000+ students)",
    ]
)

# ─── SLIDE 6: PRIORITIES WE ADDRESS ───
slide = prs.slides.add_slide(L_BULLET1)
add_bullets(slide,
    "Grant Priorities Addressed",
    "AP1 (Evidence) + AP3 (AI) + Competitive Preference",
    [
        "AP1: Early Phase - Demonstrates a Rationale",
        "   Logic model grounded in Kuh's High-Impact Practices",
        "AP3: Advancing AI to Support Student Success",
        "   EMMA uses Gemini AI for personalized coaching",
        "AP6: College-to-Career Pathways (secondary)",
        "   EMMA's architecture IS a credential map",
        "Competitive Preference: State-level partnership",
        "   UNC System Office letter = 10 bonus points",
    ]
)

# ─── SLIDE 7: PROPOSED TEAM ───
slide = prs.slides.add_slide(L_BULLET2)
add_bullets(slide,
    "Proposed Team",
    "Four Co-PIs covering every dimension",
    [
        "Dr. Nakeshia Williams - Institutional Lead",
        "   VP Undergraduate Education & Student Success",
        "Dr. Antoine Alston - College Lead",
        "   Associate Dean, CAES (25-year tenure)",
        "Charlie Hopper - Experiential Learning",
        "   Director, Small Farm Research & Innovation Center",
        "Chris Harrison - Technology",
        "   EMMA platform architect, Think! Design & Planning",
    ]
)

# ─── SLIDE 8: BUDGET OVERVIEW ───
slide = prs.slides.add_slide(L_BULLET3)
add_bullets(slide,
    "Budget Summary",
    "$3.75M over 48 months",
    [
        "Personnel (PIs, Director, coaches, GAs): $1.31M (35%)",
        "Technology Development & Infrastructure: $750K (20%)",
        "Independent Evaluation: $750K (20%)",
        "Student Support Services: $562K (15%)",
        "Other Direct Costs: $375K (10%)",
        "10% match: $375K (in-kind: PI time, facilities, EMMA IP)",
    ]
)

# ─── SLIDE 9: TWO-GRANT STRATEGY ───
slide = prs.slides.add_slide(L_BULLET4)
add_bullets(slide,
    "Two-Grant Strategy",
    "Combined: $4.25M for student success",
    [
        "NIFA Equipment Grants: up to $500K (due June 25)",
        "   Martin Building CoLab hardware & instruments",
        "   No match required",
        "",
        "FIPSE PSSG: up to $3.75M (due June 29)",
        "   EMMA + ISLA software trial & evaluation",
        "   10% match (in-kind)",
        "",
        "Hardware + Software = complete ecosystem",
    ]
)

# ─── SLIDE 10: NEXT STEPS ───
slide = prs.slides.add_slide(L_BULLET1)
add_bullets(slide,
    "What We Need This Week",
    "13 days to submission",
    [
        "Co-PI commitments from Dr. Williams & Dr. Alston",
        "Sponsored Programs notification of intent",
        "UNC System Office letter of support (10 points)",
        "Institutional retention/completion data",
        "Identify independent evaluator",
        "Chris & Charlie handle all drafting",
        "Internal review by June 26-27",
        "Submit June 29 by 11:59 PM ET",
    ]
)

# ─── CLOSING ───
slide = prs.slides.add_slide(L_TRANSITION)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 10:
        ph.text = "Questions & Discussion"

# Save
prs.save(OUTPUT)
print(f"Pitch deck saved: {OUTPUT}")
