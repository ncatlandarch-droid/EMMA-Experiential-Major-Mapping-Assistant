"""
Generate EMMA presentation using the REAL NC A&T NCATSTAT_Template.pptx
Preserves all footers, backgrounds, logos, and branding from the template.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import os

TEMPLATE = r'C:\Users\Chris\Downloads\NCATSTAT_Template.pptx'
OUTPUT   = r'C:\Users\Chris\Downloads\EMMA_Journey_Mapping_Presentation.pptx'

prs = Presentation(TEMPLATE)

# Delete all the example slides (they're blank placeholders)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Layout references
L_TITLE      = prs.slide_layouts[0]   # Title Slide
L_TRANSITION = prs.slide_layouts[2]   # Transition Slide
L_EMPTY      = prs.slide_layouts[4]   # Empty Content Slide
L_BULLET1    = prs.slide_layouts[5]   # Content Slide with Bullet Color 1
L_BULLET2    = prs.slide_layouts[6]   # Content Slide with Bullet Color 2
L_BULLET3    = prs.slide_layouts[7]   # Content Slide with Bullet Color 3
L_BULLET4    = prs.slide_layouts[8]   # Content Slide with Bullet Color 4

# ═══════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════
slide = prs.slides.add_slide(L_TITLE)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:  # Title
        ph.text = "EMMA\nStudent Experiential Journey Mapping"
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(40)
                run.font.bold = True

# ═══════════════════════════════════════
# SLIDE 2 — THE CHALLENGE (Transition)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(L_TRANSITION)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 10:
        ph.text = "Siloed Data. Fragmented Tracking.\nNo Unified View."
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(36)
                run.font.bold = True

# ═══════════════════════════════════════
# SLIDE 3 — EMMA CAPABILITY (Bullet 1)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(L_BULLET1)
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:  # Title
        ph.text = "EMMA — How It Works"
    elif idx == 10:  # Subtitle area
        ph.text = "Interactive 4-Year Journey Mapping"
    elif idx == 1:  # Body
        tf = ph.text_frame
        tf.clear()
        bullets = [
            "Student selects their college and program — 16 CAES programs loaded",
            "Sees a personalized 4-year journey map (Explore → Engage → Develop → Launch)",
            "Checks milestones as completed — each card shows course, credits, skills",
            "Progress bar animates in real time — per-category and overall",
            "Emma AI avatar speaks encouragement via Gemini Neural TTS",
            "Career strip auto-populates with live BLS salary and job outlook",
            "Exports professional PDF for advising, accreditation, employers",
        ]
        for i, b in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = b
            else:
                p = tf.add_paragraph()
                p.text = b
            para = tf.paragraphs[i]
            for run in para.runs:
                run.font.size = Pt(14)

# ═══════════════════════════════════════
# SLIDE 4 — LIVE CAREER DATA (Bullet 2)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(L_BULLET2)
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:
        ph.text = "Live Career Data — BLS Integration"
    elif idx == 10:
        ph.text = "Every Program Connects to Real Career Outcomes"
    elif idx == 1:
        tf = ph.text_frame
        tf.clear()
        bullets = [
            "16 CAES programs loaded with BLS occupation codes",
            "Landscape Architecture (17-1012) — $73,210 median salary",
            "Bioprocess Engineering (17-2031) — $99,550",
            "Natural Resources Engineering (17-2081) — $96,820",
            "Food Science (19-1012) — $80,600",
            "For each program: median salary, growth rate, top employers, certifications",
            "Data fetches LIVE from BLS.gov API — always current",
        ]
        for i, b in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = b
            else:
                p = tf.add_paragraph()
                p.text = b
            for run in tf.paragraphs[i].runs:
                run.font.size = Pt(14)

# ═══════════════════════════════════════
# SLIDE 5 — AI COACHING (Bullet 3)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(L_BULLET3)
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:
        ph.text = "AI Voice Coaching — Emma Speaks"
    elif idx == 10:
        ph.text = "Personalized Guidance Through Gemini Neural TTS"
    elif idx == 1:
        tf = ph.text_frame
        tf.clear()
        bullets = [
            "Emma is not just a dashboard — she is a coach",
            "27 pre-recorded coaching audio files per program",
            "Context-aware: knows your year, progress %, and category gaps",
            "Celebrates milestones and suggests next steps",
            "Animated avatar ring shows when Emma is speaking",
            "Four Experiential Categories (the framework):",
            "   Purpose & Self-Discovery — research, studios, capstone",
            "   Communities & Service — service-learning, community engagement",
            "   Local & Global — field work, SFRIC projects, study abroad",
            "   Professional Identity — internships, licensure prep, portfolio",
        ]
        for i, b in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = b
            else:
                p = tf.add_paragraph()
                p.text = b
            for run in tf.paragraphs[i].runs:
                run.font.size = Pt(13)

# ═══════════════════════════════════════
# SLIDE 6 — BSLA PILOT (Bullet 4)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(L_BULLET4)
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:
        ph.text = "BSLA Pilot — Deep Curriculum Integration"
    elif idx == 10:
        ph.text = "The Fully Realized Proof of Concept"
    elif idx == 1:
        tf = ph.text_frame
        tf.clear()
        bullets = [
            "24 courses across 4 years, each with credits, semester, instructor",
            "LAAB accreditation standards mapped (3.A, 3.B.1–3.B.10)",
            "Learning objectives by domain + skills inventory",
            "4 SFRIC Field Projects:",
            "   Inspiration Courtyard — CH Moore Ag Research Station",
            "   Holland Bowl — NC A&T University Farm",
            "   Woodland Restoration — Woodland Edge",
            "   Small Area Plan — SFRIC Community Zone",
            "6 ASLA Student Award categories with deliverable checklists",
            "This level of detail can be replicated for ANY program",
        ]
        for i, b in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = b
            else:
                p = tf.add_paragraph()
                p.text = b
            for run in tf.paragraphs[i].runs:
                run.font.size = Pt(13)

# ═══════════════════════════════════════
# SLIDE 7 — ECOSYSTEM (Bullet 1)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(L_BULLET1)
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:
        ph.text = "The Connected Ecosystem"
    elif idx == 10:
        ph.text = "EMMA Doesn't Stand Alone — 4 Platforms, 1 Vision"
    elif idx == 1:
        tf = ph.text_frame
        tf.clear()
        bullets = [
            "EMMA — Experiential Major Mapping Assistant (thinkemma.app)",
            "   Interactive journey mapping + AI coaching + career data",
            "AVA — Adaptive Visualization Assistant (ava3-digital-twin.netlify.app)",
            "   CesiumJS 3D digital twin + SITES v2 scoring (250 pts)",
            "GRANT — Grant Research & Award Navigation (thinkgranted.app)",
            "   AI-guided funding pipeline — NOI wizard to award",
            "ISLA — Interactive Study & Licensure Assistant (larelab.app)",
            "   Adaptive exam prep + Coach Perry AI + spaced repetition",
            "",
            "All platforms share Firebase Auth + Firestore + design language",
        ]
        for i, b in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = b
            else:
                p = tf.add_paragraph()
                p.text = b
            for run in tf.paragraphs[i].runs:
                run.font.size = Pt(13)

# ═══════════════════════════════════════
# SLIDE 8 — CANVAS INTEGRATION (Bullet 2)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(L_BULLET2)
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:
        ph.text = "Canvas LMS Integration — Ready When You Are"
    elif idx == 10:
        ph.text = "LTI 1.3 — Seamless Canvas Connection"
    elif idx == 1:
        tf = ph.text_frame
        tf.clear()
        bullets = [
            "Level 1: Embedded Link (1 week) — CAN BE DEPLOYED IMMEDIATELY",
            "   EMMA opens inside Canvas as a module link. No coding needed.",
            "Level 2: LTI 1.3 Basic (2–3 weeks)",
            "   Full SSO with NC A&T identity. Auto-detects program.",
            "Level 3: LTI 1.3 + Grade Passback (4–6 weeks)",
            "   Milestone completions push to Canvas gradebook.",
            "Level 4: Deep Integration (8–12 weeks)",
            "   Deep linking, roster sync, advisor analytics dashboard.",
            "",
            "FERPA Compliant • We need: Canvas LTI Developer Key from NC A&T IT",
        ]
        for i, b in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = b
            else:
                p = tf.add_paragraph()
                p.text = b
            for run in tf.paragraphs[i].runs:
                run.font.size = Pt(13)

# ═══════════════════════════════════════
# SLIDE 9 — FEEDBACK ASK (Bullet 3)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(L_BULLET3)
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:
        ph.text = "We Need Your Input"
    elif idx == 10:
        ph.text = "Collaborative Design Phase"
    elif idx == 1:
        tf = ph.text_frame
        tf.clear()
        bullets = [
            "Workflow Fit:",
            "   Does the 4-category framework align with your program?",
            "   What milestones matter most for YOUR students?",
            "Usability:",
            "   Would students use this voluntarily — or need it assigned?",
            "   Should AI coaching be opt-in or always-on?",
            "Strategic Impact:",
            "   How does this support your accreditation self-study?",
            "   Would Canvas integration change your adoption calculus?",
            "",
            "Send us your milestone data. We will load your program.",
        ]
        for i, b in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = b
            else:
                p = tf.add_paragraph()
                p.text = b
            for run in tf.paragraphs[i].runs:
                run.font.size = Pt(13)

# ═══════════════════════════════════════
# SLIDE 10 — Q&A (Transition)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(L_TRANSITION)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 10:
        ph.text = "Questions & Live Demo\n\nthinkemma.app\n\nW. Christopher Harrison | wcharris@ncat.edu\n© 2026 Think! Design and Planning, LLC"
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(24)

# ═══════════════════════════════════════
# SAVE
# ═══════════════════════════════════════
prs.save(OUTPUT)
print(f'✅ Saved to: {OUTPUT}')
print(f'   {len(prs.slides)} slides using REAL NC A&T template')
